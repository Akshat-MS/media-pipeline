#!/usr/bin/env python3
"""
mpk — Media Pipeline Kit

Command-line utilities for the Lecture Alive AI workflow. Deterministic work
lives here; judgement lives in the layer prompts.

    mpk --help
    mpk deck --help
    mpk deck extract --help
"""
from __future__ import annotations

import argparse
import base64
import copy
import json
import re
import shutil
import tempfile
import subprocess
import sys
from pathlib import Path

__version__ = "0.1.0"

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
EMU_PER_INCH = 914400


# ════════════════════════════════════════════════════════════ helpers
def die(msg: str, code: int = 1):
    print(f"mpk: error: {msg}", file=sys.stderr)
    sys.exit(code)


def need(mod: str):
    try:
        return __import__(mod)
    except ImportError:
        die(f"{mod} is required — pip install {mod.replace('_','-')}")


def have(cmd: str) -> bool:
    return shutil.which(cmd) is not None


def safe(fn, default=None):
    """python-pptx raises rather than returning None for several properties
    (auto_shape_type on a non-autoshape, colour on an inherited fill). Probing
    is normal here, so an exception is a fact about the shape, not an error."""
    try:
        return fn()
    except Exception:
        return default


def emit(data, out: str | None):
    text = json.dumps(data, indent=2, ensure_ascii=False)
    if out:
        Path(out).write_text(text, encoding="utf-8")
        print(f"wrote {out}  ({len(text):,} bytes)", file=sys.stderr)
    else:
        print(text)


def slide_id(deck_id: str, n: int) -> str:
    return f"{deck_id}_s{n:02d}"


def open_deck(path: str):
    """Open a .pptx. A file that is not a PowerPoint package is a normal
    mistake, not a crash — python-pptx raises PackageNotFoundError, which says
    nothing useful about what went wrong."""
    need("pptx")
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError
    try:
        return Presentation(path)
    except PackageNotFoundError:
        die(f"not a PowerPoint file: {path}\n"
            f"       mpk deck commands need a .pptx. If this is a .ppt, open it "
            f"in PowerPoint or LibreOffice and save as .pptx.")
    except Exception as exc:
        die(f"could not open {path}: {type(exc).__name__}: {exc}")


# ════════════════════════════════════════════════════════════ deck extract
def _lines_of(tf) -> list[str]:
    """Visual lines. Splits paragraphs on <a:br/> soft breaks, which PowerPoint
    uses freely and which paragraph.text silently joins with \\x0b — merging
    what a viewer sees as separate lines. Layer 8 reveals visual lines, so this
    is the list that matters."""
    out = []
    for para in tf.paragraphs:
        for piece in para.text.split("\x0b"):
            out.append(piece)
    return out


def _colour(obj, kind: str) -> dict | None:
    """Best-effort colour read. Records the token as well as the resolved RGB,
    because a theme reference and a literal hex are different facts."""
    try:
        src = getattr(obj, kind)
        col = src.fore_color if kind == "fill" else src.color
        t = str(col.type).split(".")[-1].split(" ")[0] if col.type is not None else None
        d = {"type": t}
        try:
            d["rgb"] = str(col.rgb)
        except Exception:
            pass
        try:
            d["theme_color"] = str(col.theme_color).split(".")[-1].split(" ")[0]
        except Exception:
            pass
        return d if len(d) > 1 else None
    except Exception:
        return None


def _connector_facts(sh) -> dict:
    """Direction is the highest-error field in this layer, so read it from the
    XML rather than inferring it: a:stCxn / a:endCxn give the connected shapes,
    a:tailEnd / a:headEnd give which end carries the arrowhead."""
    el = sh._element
    out = {"start_cxn": None, "end_cxn": None,
           "head_end": None, "tail_end": None, "direction_verified": False}

    nv = el.find(f"{P}nvCxnSpPr/{P}cNvCxnSpPr")
    if nv is not None:
        for tag, key in (("stCxn", "start_cxn"), ("endCxn", "end_cxn")):
            n = nv.find(f"{A}{tag}")
            if n is not None:
                out[key] = {"shape_id": int(n.get("id")), "site_idx": int(n.get("idx"))}

    ln = el.find(f"{P}spPr/{A}ln")
    if ln is not None:
        for tag, key in (("headEnd", "head_end"), ("tailEnd", "tail_end")):
            n = ln.find(f"{A}{tag}")
            if n is not None:
                out[key] = n.get("type")

    arrow = (out["tail_end"] and out["tail_end"] != "none") or \
            (out["head_end"] and out["head_end"] != "none")
    if out["start_cxn"] and out["end_cxn"] and arrow:
        out["direction_verified"] = True
        out["direction_basis"] = (
            "a:stCxn/a:endCxn give both connected shapes; "
            f"a:tailEnd type={out['tail_end']!r} head={out['head_end']!r} "
            "identifies the arrowhead end")
    else:
        missing = []
        if not out["start_cxn"] or not out["end_cxn"]:
            missing.append("connection sites not declared in XML")
        if not arrow:
            missing.append("no arrowhead on either end")
        out["direction_basis"] = "NOT VERIFIED — " + "; ".join(missing)
    return out


def _table_of(sh) -> dict | None:
    if not getattr(sh, "has_table", False):
        return None
    t = sh.table
    return {
        "rows": len(t.rows), "cols": len(t.columns),
        "cells": [[c.text for c in row.cells] for row in t.rows],
    }


def _walk(shapes, sid: str, parent=None, group_off=(0, 0), depth=0) -> list[dict]:
    out = []
    for sh in shapes:
        eid = f"{sid}_{sh.shape_id}"
        raw = {"x": sh.left, "y": sh.top, "cx": sh.width, "cy": sh.height}
        absbox = dict(raw)
        if parent is not None:
            absbox = {"x": (sh.left or 0) + group_off[0], "y": (sh.top or 0) + group_off[1],
                      "cx": sh.width, "cy": sh.height}

        st = str(safe(lambda: sh.shape_type) or "")
        is_group = "GROUP" in st.upper()
        is_conn = sh._element.tag.endswith("}cxnSp")
        is_pic = "PICTURE" in st.upper()
        is_tbl = getattr(sh, "has_table", False)

        a = {
            "element_id": eid,
            "ooxml_shape_id": sh.shape_id,
            "ooxml_name": sh.name,
            "type": ("connector" if is_conn else "group" if is_group
                     else "table" if is_tbl else "picture" if is_pic
                     else "text" if safe(lambda: bool(sh.text_frame.text.strip()), False)
                     else "shape"),
            "shape_type": st,
            "shape_kind": safe(lambda: str(sh.auto_shape_type).split(" ")[0]),
            "parent_id": parent,
            "bounding_box": absbox,
            "properties": {},
            "flags": [],
        }
        if parent is not None:
            a["bounding_box_raw"] = raw
            a["flags"].append({
                "code": "group_relative_geometry", "severity": "info",
                "note": "inside a group — raw offsets are relative to the group transform; "
                        "bounding_box is the resolved absolute box",
                "needs_eye_check": False})

        if safe(lambda: sh.rotation):
            a["bounding_box"]["rot"] = sh.rotation
            a["flags"].append({"code": "rotated", "severity": "info",
                               "note": f"rotation {sh.rotation}° — resolve before "
                                       "concluding anything about geometry",
                               "needs_eye_check": False})

        p = a["properties"]
        if safe(lambda: sh.has_text_frame, False):
            tf = sh.text_frame
            p["text"] = tf.text
            p["paragraphs"] = [{"level": q.level, "text": q.text} for q in tf.paragraphs]
            p["lines"] = _lines_of(tf)
            if len(p["lines"]) > len(p["paragraphs"]):
                a["flags"].append({
                    "code": "soft_line_breaks", "severity": "info",
                    "note": f"{len(p['lines'])} visual lines across "
                            f"{len(p['paragraphs'])} paragraphs — <a:br/> soft breaks "
                            "present; use lines[] for per-line reveal",
                    "needs_eye_check": False})
            szs = {r.font.size.pt for q in tf.paragraphs for r in q.runs
                   if r.font.size is not None}
            if szs:
                p["font_sizes_pt"] = sorted(szs)
            if not tf.text.strip():
                p["empty"] = True
                a["flags"].append({"code": "empty_text_box", "severity": "info",
                                   "note": "no text runs — contributes no content",
                                   "needs_eye_check": False})
            if re.search(r"</?a:fld", sh._element.xml):
                p["is_field"] = True
                a["flags"].append({"code": "rendered_field", "severity": "info",
                                   "note": "contains a field (e.g. slide number) — "
                                           "the value is rendered, not literal",
                                   "needs_eye_check": False})

        for k, kind in (("fill", "fill"), ("line", "line")):
            c = _colour(sh, kind)
            if c:
                p[f"{k}_colour"] = c

        if is_pic:
            p["descr"] = safe(lambda: sh._element.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}"
                "..").get("descr")) or sh.name
            p["image_name"] = safe(lambda: sh.image.filename)
            p["image_ext"] = safe(lambda: sh.image.ext)
            p["raster_content_unread"] = True
            a["flags"].append({"code": "raster_content_unread", "severity": "warn",
                               "note": "opaque raster — the XML cannot say what is "
                                       "drawn inside it",
                               "needs_eye_check": True})

        if is_conn:
            cf = _connector_facts(sh)
            a["endpoints"] = {
                "start": ({"anchor_shape_id": cf["start_cxn"]["shape_id"],
                           "anchor_element_id": f"{sid}_{cf['start_cxn']['shape_id']}",
                           "site_idx": cf["start_cxn"]["site_idx"], "inferred": False}
                          if cf["start_cxn"] else None),
                "end": ({"anchor_shape_id": cf["end_cxn"]["shape_id"],
                         "anchor_element_id": f"{sid}_{cf['end_cxn']['shape_id']}",
                         "site_idx": cf["end_cxn"]["site_idx"], "inferred": False}
                        if cf["end_cxn"] else None)}
            p.update({k: cf[k] for k in
                      ("head_end", "tail_end", "direction_verified", "direction_basis")})
            if not cf["direction_verified"]:
                a["flags"].append({"code": "direction_unverified", "severity": "error",
                                   "note": cf["direction_basis"],
                                   "needs_eye_check": True})

        if is_tbl:
            p["table"] = _table_of(sh)

        out.append(a)

        if is_group:
            off = (absbox["x"], absbox["y"])
            out.extend(_walk(sh.shapes, sid, parent=eid, group_off=off, depth=depth + 1))
    return out


def cmd_deck_extract(args):
    prs = open_deck(args.file)
    deck_id = args.deck_id or Path(args.file).stem.lower()[:16]

    slides, coverage = [], []
    for n, s in enumerate(prs.slides, 1):
        if args.slide and n != args.slide:
            continue
        sid = slide_id(deck_id, n)
        assets = _walk(s.shapes, sid)
        ids = sorted(a["ooxml_shape_id"] for a in assets)
        gaps = [i for i in range(min(ids), max(ids) + 1) if i not in ids] if ids else []
        coverage.append({"slide_id": sid, "captured_ids": ids, "gaps": gaps,
                         "reading": ("no shapes at these ids in spTree — consistent with "
                                     "deleted shapes or non-spTree parts"
                                     if gaps else "contiguous, no gaps")})
        notes = None
        if s.has_notes_slide and s.notes_slide.notes_text_frame is not None:
            t = s.notes_slide.notes_text_frame.text.strip()
            notes = t or None
        slides.append({"slide_id": sid, "slide_number": n,
                       "notes": notes, "assets": assets})

    out = {
        "deck_id": deck_id,
        "metadata": {
            "source_file": Path(args.file).name,
            "slide_count": len(prs.slides),
            "slides_extracted": len(slides),
            "canvas_dimensions": {"width": prs.slide_width, "height": prs.slide_height,
                                  "units": "EMU",
                                  "inches": [round(prs.slide_width / EMU_PER_INCH, 3),
                                             round(prs.slide_height / EMU_PER_INCH, 3)],
                                  "aspect": round(prs.slide_width / prs.slide_height, 4)},
            "extraction_path": "mpk_deck_extract",
            "mpk_version": __version__,
        },
        "id_coverage": coverage,
        "slides": slides,
    }
    if args.expect and args.expect != len(prs.slides):
        out["metadata"]["slide_count_mismatch"] = {
            "expected": args.expect, "found": len(prs.slides)}
        print(f"mpk: WARNING expected {args.expect} slides, found {len(prs.slides)}",
              file=sys.stderr)
    emit(out, args.out)


def cmd_deck_info(args):
    prs = open_deck(args.file)
    w, h = prs.slide_width, prs.slide_height
    print(f"{Path(args.file).name}")
    print(f"  slides : {len(prs.slides)}")
    print(f"  canvas : {w} x {h} EMU  "
          f"({w/EMU_PER_INCH:.2f} x {h/EMU_PER_INCH:.2f} in, "
          f"aspect {w/h:.4f})")
    for n, s in enumerate(prs.slides, 1):
        kinds: dict[str, int] = {}
        for sh in s.shapes:
            k = str(safe(lambda: sh.shape_type) or "?").split(" ")[0]
            kinds[k] = kinds.get(k, 0) + 1
        title = next((sh.text_frame.text.split("\n")[0][:48]
                      for sh in s.shapes
                      if safe(lambda: bool(sh.text_frame.text.strip()), False)
                      and not sh.text_frame.text.strip().isdigit()), "")
        bits = " ".join(f"{k}×{v}" for k, v in sorted(kinds.items()))
        print(f"  s{n:02d}  {len(s.shapes):2d} shapes  {bits:<48} {title}")


# ════════════════════════════════════════════════════════════ deck normalize
def cmd_deck_normalize(args):
    src = open_deck(args.file)
    if args.like:
        ref = open_deck(args.like)
        tw, th = ref.slide_width, ref.slide_height
    else:
        tw, th = args.width, args.height
    if not tw or not th:
        die("give --like <deck.pptx> or both --width and --height")

    sw, sh_ = src.slide_width, src.slide_height
    fx, fy = tw / sw, th / sh_
    if abs((sw / sh_) - (tw / th)) > 0.01:
        print(f"mpk: WARNING aspect differs ({sw/sh_:.4f} vs {tw/th:.4f}) — "
              "scaling will distort", file=sys.stderr)

    def scale(shapes):
        for s in shapes:
            for attr, f in (("left", fx), ("top", fy), ("width", fx), ("height", fy)):
                v = getattr(s, attr)
                if v is not None:
                    setattr(s, attr, int(round(v * f)))
            if str(safe(lambda: s.shape_type) or "").upper().startswith("GROUP"):
                scale(s.shapes)
            if safe(lambda: s.has_text_frame, False):
                for para in s.text_frame.paragraphs:
                    for r in para.runs:
                        if r.font.size is not None:
                            r.font.size = int(round(r.font.size * min(fx, fy)))

    for s in src.slides:
        scale(s.shapes)
    src.slide_width, src.slide_height = int(tw), int(th)
    src.save(args.out)
    print(f"normalized {sw}x{sh_} -> {int(tw)}x{int(th)}  "
          f"(x{fx:.4f}, y{fy:.4f})  -> {args.out}", file=sys.stderr)


# ════════════════════════════════════════════════════════════ deck merge
def cmd_deck_merge(args):
    base = open_deck(args.into)
    add = open_deck(args.file)

    if (base.slide_width, base.slide_height) != (add.slide_width, add.slide_height):
        die(f"canvas mismatch: {args.into} is {base.slide_width}x{base.slide_height}, "
            f"{args.file} is {add.slide_width}x{add.slide_height}. "
            f"Run 'mpk deck normalize {args.file} --like {args.into}' first — merging "
            f"decks of different sizes puts every coordinate in the wrong frame.")

    layout = base.slide_layouts[args.layout] if args.layout is not None \
        else base.slides[0].slide_layout

    inserted = []
    for s in add.slides:
        new = base.slides.add_slide(layout)
        for shp in list(new.shapes):          # start from a clean slide
            shp._element.getparent().remove(shp._element)
        for shp in s.shapes:
            new.shapes._spTree.append(copy.deepcopy(shp._element))
        inserted.append(new)

    if args.at:
        sldIdLst = base.slides._sldIdLst
        ids = list(sldIdLst)
        for k, new in enumerate(inserted):
            el = ids[len(ids) - len(inserted) + k]
            sldIdLst.remove(el)
            sldIdLst.insert(args.at - 1 + k, el)

    base.save(args.out)
    print(f"merged {len(inserted)} slide(s) from {Path(args.file).name} into "
          f"{Path(args.into).name}"
          + (f" at position {args.at}" if args.at else " (appended)")
          + f"  -> {args.out}", file=sys.stderr)
    print("mpk: NOTE images and theme parts are not re-linked by a raw spTree copy — "
          "run 'mpk deck render' and check the result before trusting it.",
          file=sys.stderr)


# ════════════════════════════════════════════════════════════ deck render
def cmd_deck_render(args):
    soffice = next((c for c in ("soffice", "libreoffice") if have(c)), None)
    if not soffice:
        die("LibreOffice not found — install it, or skip rendering")
    outdir = Path(args.outdir)
    outdir.mkdir(parents=True, exist_ok=True)
    pdf = outdir / (Path(args.file).stem + ".pdf")
    subprocess.run([soffice, "--headless", "--convert-to", "pdf",
                    "--outdir", str(outdir), args.file], check=True,
                   stdout=subprocess.DEVNULL)
    if not have("pdftoppm"):
        print(f"wrote {pdf} — install poppler-utils (pdftoppm) for PNGs",
              file=sys.stderr)
        return
    subprocess.run(["pdftoppm", "-png", "-r", str(args.dpi),
                    str(pdf), str(outdir / Path(args.file).stem)], check=True)
    pngs = sorted(outdir.glob(f"{Path(args.file).stem}*.png"))
    print(f"rendered {len(pngs)} page(s) at {args.dpi} dpi -> {outdir}", file=sys.stderr)
    for p in pngs:
        print(p)


# ════════════════════════════════════════════════════════════ review build
TEMPLATE_DIR = Path(__file__).parent / "templates"
DEFAULT_TEMPLATE = "slide-review"
_DATA_BLOCK = re.compile(
    r'(<script type="application/json" id="manifest">\n).*?(\n</script>)', re.DOTALL)


def resolve_template(name: str) -> Path:
    """Accept a bare template name or an explicit path.

    Bare names resolve against tools/templates/ so callers say
    `--template slide-review` rather than carrying a path around. Later layers
    add their own pages here — a sequence player for Layer 8, a transcript
    table for Layer 4 — and each is a template, not a new renderer."""
    p = Path(name)
    if p.suffix and p.exists():
        return p
    cand = TEMPLATE_DIR / (name if name.endswith(".html") else f"{name}.html")
    if cand.exists():
        return cand
    have_names = ", ".join(sorted(t.stem for t in TEMPLATE_DIR.glob("*.html"))) or "none"
    die(f"no template {name!r} — available: {have_names}")


def cmd_review_templates(args):
    if not TEMPLATE_DIR.is_dir():
        die(f"no templates directory at {TEMPLATE_DIR}")
    rows = sorted(TEMPLATE_DIR.glob("*.html"))
    if not rows:
        print("no templates found")
        return
    print(f"templates in {TEMPLATE_DIR}:")
    for t in rows:
        m = re.search(r"<title>(.*?)</title>", t.read_text(encoding="utf-8"))
        mark = "  (default)" if t.stem == DEFAULT_TEMPLATE else ""
        print(f"  {t.stem:<20} {m.group(1) if m else '':<40}{mark}")


def cmd_review_build(args):
    """Inject the manifest into a review template.

    The template owns the page — its embedded JavaScript renders every row,
    flag, prior and connector sentence from the manifest at load time. This
    command only substitutes the data block, so there is exactly one
    implementation of the page, shared by the tool and by the prompt."""
    tpl_path = resolve_template(args.template or DEFAULT_TEMPLATE)
    tpl = tpl_path.read_text(encoding="utf-8")
    raw = Path(args.file).read_text(encoding="utf-8")
    try:
        json.loads(raw)                  # fail loudly, and readably
    except json.JSONDecodeError as exc:
        die(f"{args.file} is not valid JSON: {exc}")
    html, n = _DATA_BLOCK.subn(
        lambda m: m.group(1) + raw.strip() + m.group(2), tpl)
    if n != 1:
        die(f'expected exactly one manifest block in {tpl_path}, found {n}')
    Path(args.out).write_text(html, encoding="utf-8")
    print(f"wrote {args.out}  ({len(html):,} bytes) — {tpl_path.name}, "
          f"rendering from the embedded manifest", file=sys.stderr)



# ════════════════════════════════════════════════════════════ transcript
# Whisper's own vocabulary bias. Notation is where transcription fails, and
# biasing the recogniser up front beats detecting the error afterwards —
# so this is the cheapest place to act on OBS-021.
DEFAULT_VOCAB = ("P1, P2, P3, Pn, Pi, R1, R2, R3, Rn, Rj, semaphore, mutex, "
                 "bounded buffer, critical section, deadlock, BUFSIZE")

# Symbolic notation, deck-independent: subscripts, P1/Rj forms, set braces,
# relational and arrow operators.
NOTATION_RE = re.compile(
    r"[\u2080-\u2089\u1d62\u2c7c\u2099]|\b[PRp]\s?[0-9ijn]\b|[{}]|\u2264|\u2192|-->|\bP\s?sub\b")


def terms_from_manifest(path: str) -> list[str]:
    """Pull the deck's own vocabulary out of a Layer 3 manifest.

    A hardcoded pattern only knows one deck's notation. The V017 run proved
    it: the regex hunts P1/Rj — deadlock vocabulary — and flagged ZERO
    segments in a transcript saying "semaphore" 8 times, "mutex" 7 and
    "buff size" once. The deck already knows its own terms, so read them
    from there instead of maintaining a regex per topic."""
    try:
        d = json.loads(Path(path).read_text(encoding="utf-8"))
    except Exception as exc:
        die(f"could not read manifest {path}: {exc}")
    d = d.get("payload", d)
    terms: set[str] = set()

    for e in (d.get("deck") or {}).get("entity_inventory") or []:
        if e.get("label"):
            terms.add(str(e["label"]).strip())

    for sl in d.get("slides") or [d]:
        for a in sl.get("assets") or []:
            p = a.get("properties") or {}
            for line in (p.get("lines") or []):
                for tok in re.findall(r"[A-Za-z_][A-Za-z0-9_]{2,}", line or ""):
                    if tok.isupper() or re.search(r"[a-z][A-Z]|_", tok):
                        terms.add(tok)          # BUFSIZE, next_produced, camelCase
    return sorted(t for t in terms if len(t) > 1)


def _fuzzy_term(t: str) -> str:
    """A pattern that survives the ways a recogniser mangles an identifier.

    V017 proved the need: the deck says BUFSIZE, the recogniser wrote
    "buff size" — split in a different place than expected, with a doubled
    letter. So allow a separator between ANY two characters, and allow any
    character to repeat. Only for terms long enough that this stays specific."""
    if len(t) < 5 or not t.isalpha():
        return re.escape(t)
    return r"[\s\-]?".join(re.escape(c) + "+" for c in t)


def build_notation_matcher(terms: list[str]):
    """Match symbolic notation, or any term we have declared we care about.

    Two sources, both needed. The deck supplies identifiers (BUFSIZE,
    next_produced); the vocabulary bias supplies the domain words (semaphore,
    mutex, wait, signal) that never appear as slide text but carry the meaning.
    Whatever we told the recogniser to listen for is, by definition, what we
    need it to have got right."""
    if not terms:
        return lambda text: bool(NOTATION_RE.search(text or ""))
    alts = sorted({_fuzzy_term(t) for t in terms if len(t) > 2},
                  key=len, reverse=True)
    term_re = re.compile(r"(?<![A-Za-z])(?:" + "|".join(alts) + r")(?![A-Za-z])",
                         re.IGNORECASE)
    return lambda text: bool(NOTATION_RE.search(text or "")
                             or term_re.search(text or ""))


def _split_vocab(v: str | None) -> list[str]:
    """Vocabulary bias is a comma-separated string; multi-word phrases count."""
    return [t.strip() for t in (v or "").split(",") if t.strip()]


def cmd_transcript_build(args):
    """Transcribe and force-align. Emits RAW ASR output — the verification
    prompt turns this into the final transcript, exactly as `deck extract`
    feeds the Layer 3 prompt."""
    need("faster_whisper")
    from faster_whisper import WhisperModel

    vocab = args.vocab
    if args.vocab_file:
        vocab = Path(args.vocab_file).read_text(encoding="utf-8").strip()
    deck_terms: list[str] = []
    if args.terms_from:
        deck_terms = terms_from_manifest(args.terms_from)
        print(f"mpk: {len(deck_terms)} deck term(s) from "
              f"{Path(args.terms_from).name} — notation detection will use them",
              file=sys.stderr)
        if vocab is None:
            vocab = ", ".join(deck_terms)      # the deck is the better bias too
    if vocab is None:
        vocab = DEFAULT_VOCAB
    # Everything we asked the recogniser to listen for is something we need it
    # to have got right — so the bias list is also the flag list.
    critical = sorted(set(deck_terms) | set(_split_vocab(vocab)))
    is_notation = build_notation_matcher(critical)
    print(f"mpk: {len(critical)} critical term(s) drive notation flagging",
          file=sys.stderr)

    print(f"mpk: loading model {args.model!r} ({args.compute}) …", file=sys.stderr)
    try:
        model = WhisperModel(args.model, device=args.device,
                             compute_type=args.compute,
                             download_root=args.model_dir or None)
    except Exception as exc:
        die(f"could not load model {args.model!r}: {type(exc).__name__}: {exc}\n"
            f"       Models download from Hugging Face on first use. If the network\n"
            f"       blocks that, fetch the model once elsewhere and point at it with\n"
            f"       --model-dir, or pass a local path to --model.")

    segments, info = model.transcribe(
        args.file,
        language=args.language,
        word_timestamps=True,          # DEC-001: word timings are mandatory
        vad_filter=not args.no_vad,
        initial_prompt=vocab or None,
        beam_size=args.beam,
    )

    segs, n_words = [], 0
    for sg in segments:
        words = [{"word": w.word.strip(), "start": round(w.start, 3),
                  "end": round(w.end, 3),
                  "probability": round(w.probability, 4)}
                 for w in (sg.words or [])]
        n_words += len(words)
        low = [w for w in words if w["probability"] < args.low_conf]
        flags = []
        if not words:
            flags.append({"code": "no_word_timings", "severity": "error",
                          "note": "segment has no word-level timings; Layer 8 "
                                  "cannot bind a beat inside it",
                          "needs_eye_check": True})
        if low:
            flags.append({"code": "low_confidence_words", "severity": "warn",
                          "note": f"{len(low)} word(s) below {args.low_conf} "
                                  f"confidence: " + ", ".join(w["word"] for w in low[:6]),
                          "needs_eye_check": True})
        notation = is_notation(sg.text or "")
        if notation:
            flags.append({"code": "notation_present", "severity": "warn",
                          "note": "contains notation or a deck term — confirm "
                                  "it against the audio before trusting it "
                                  "(OBS-021)",
                          "needs_eye_check": True})
        segs.append({"segment_id": f"t_{len(segs)+1:04d}",
                     "start": round(sg.start, 3), "end": round(sg.end, 3),
                     "text": (sg.text or "").strip(),
                     "notation": notation,
                     "avg_logprob": round(sg.avg_logprob, 4),
                     "no_speech_prob": round(sg.no_speech_prob, 4),
                     "words": words, "flags": flags})

    pauses = [round(b["start"] - a["end"], 3)
              for a, b in zip(segs, segs[1:]) if b["start"] - a["end"] >= args.pause]

    out = {
        "video_id": args.video_id or Path(args.file).stem,
        "metadata": {
            "source_file": Path(args.file).name,
            "model": args.model, "device": args.device,
            "compute_type": args.compute,
            "language": info.language,
            "language_probability": round(info.language_probability, 4),
            "duration_s": round(info.duration, 3),
            "granularity": "word",
            "word_level_available": n_words > 0,
            "word_count": n_words,
            "vad_filter": not args.no_vad,
            "vocab_bias": vocab,
            "notation_terms": critical or None,
            "notation_terms_source": (args.terms_from if args.terms_from
                                      else "built-in symbolic pattern only"),
            "pause_threshold_s": args.pause,
            "pauses_over_threshold": len(pauses),
            "longest_pause_s": max(pauses) if pauses else None,
            "extraction_path": "mpk_transcript_build",
            "mpk_version": __version__,
            "note": "RAW ASR output. Verbatim — not cleaned, summarised or "
                    "reordered. The narration is fixed (project invariant); "
                    "this records when each word was said, nothing more.",
        },
        "segments": segs,
    }
    emit(out, args.out)
    print(f"mpk: {len(segs)} segments, {n_words} words, "
          f"{sum(1 for s in segs if s['notation'])} notation-bearing, "
          f"{len(pauses)} pauses over {args.pause}s", file=sys.stderr)


def _vtt(t) -> str:
    def ts(s):
        h, r = divmod(float(s), 3600); m, sec = divmod(r, 60)
        return f"{int(h):02d}:{int(m):02d}:{sec:06.3f}"
    lines = ["WEBVTT", ""]
    for i, sg in enumerate(t.get("segments", []), 1):
        lines += [str(i), f"{ts(sg['start'])} --> {ts(sg['end'])}",
                  sg.get("text", "").strip(), ""]
    return "\n".join(lines)


def _txt(t) -> str:
    """The stamped, paste-able form — a timestamp line then its text, the way a
    transcript is read and corrected by hand."""
    def mmss(s):
        m, sec = divmod(int(float(s)), 60)
        return f"{m:02d}:{sec:02d}"
    out = []
    for sg in t.get("segments", []):
        out.append(mmss(sg["start"]))
        out.append(sg.get("text", "").strip())
    return "\n".join(out) + "\n"


def _audio_data_uri(path: str) -> str:
    import base64, mimetypes
    mt = mimetypes.guess_type(path)[0] or "audio/mpeg"
    b = Path(path).read_bytes()
    print(f"mpk: embedding {len(b)/1e6:.1f} MB of audio "
          f"({len(b)*4//3/1e6:.1f} MB once base64-encoded)", file=sys.stderr)
    return f"data:{mt};base64," + base64.b64encode(b).decode("ascii")


_TRANSCRIPT_BLOCK = re.compile(
    r'(<script type="application/json" id="transcript">\n).*?(\n</script>)', re.DOTALL)
_AUDIO_BLOCK = re.compile(
    r'(<script type="text/plain" id="audio-data">).*?(</script>)', re.DOTALL)


def cmd_transcript_export(args):
    raw = Path(args.file).read_text(encoding="utf-8")
    try:
        t = json.loads(raw)
    except json.JSONDecodeError as exc:
        die(f"{args.file} is not valid JSON: {exc}")
    t = t.get("payload", t)

    fmt = args.format
    if fmt == "json":
        Path(args.out).write_text(json.dumps(t, indent=2, ensure_ascii=False),
                                  encoding="utf-8")
    elif fmt == "vtt":
        Path(args.out).write_text(_vtt(t), encoding="utf-8")
    elif fmt == "txt":
        Path(args.out).write_text(_txt(t), encoding="utf-8")
    elif fmt == "html":
        tpl = resolve_template(args.template or "transcript-review").read_text(
            encoding="utf-8")
        html, n = _TRANSCRIPT_BLOCK.subn(
            lambda m: m.group(1) + json.dumps(t, indent=1, ensure_ascii=False)
            + m.group(2), tpl)
        if n != 1:
            die("expected exactly one transcript block in the template")
        if args.audio:
            if not Path(args.audio).exists():
                die(f"audio not found: {args.audio}")
            uri = _audio_data_uri(args.audio)
            html, n = _AUDIO_BLOCK.subn(lambda m: m.group(1) + uri + m.group(2), html)
            if n != 1:
                die("expected exactly one audio block in the template")
        else:
            print("mpk: no --audio given — the page will render but timestamps "
                  "will not play", file=sys.stderr)
        Path(args.out).write_text(html, encoding="utf-8")
    else:
        die(f"unknown format {fmt!r}")
    print(f"wrote {args.out}  ({Path(args.out).stat().st_size:,} bytes, {fmt})",
          file=sys.stderr)


def cmd_transcript_reflag(args):
    """Re-apply notation detection to an existing transcript.

    Transcribing again costs minutes; re-flagging costs nothing. Use this when
    the deck terms change, or when a transcript was built before its Layer 3
    manifest existed."""
    raw = Path(args.file).read_text(encoding="utf-8")
    try:
        t = json.loads(raw)
    except json.JSONDecodeError as exc:
        die(f"{args.file} is not valid JSON: {exc}")
    t = t.get("payload", t)
    terms = sorted(set(terms_from_manifest(args.terms_from) if args.terms_from else [])
                   | set(_split_vocab(args.terms
                                      or (t.get("metadata") or {}).get("vocab_bias"))))
    is_notation = build_notation_matcher(terms)

    before = sum(1 for s_ in t.get("segments", []) if s_.get("notation"))
    changed = 0
    for s_ in t.get("segments", []):
        was, now = bool(s_.get("notation")), is_notation(s_.get("text", ""))
        if was != now:
            changed += 1
        s_["notation"] = now
        s_["flags"] = [f for f in (s_.get("flags") or [])
                       if f.get("code") != "notation_present"]
        if now:
            s_["flags"].append({
                "code": "notation_present", "severity": "warn",
                "note": "contains notation or a deck term — confirm it against "
                        "the audio before trusting it (OBS-021)",
                "needs_eye_check": True})
    t.setdefault("metadata", {})
    t["metadata"]["notation_terms"] = terms or None
    t["metadata"]["notation_terms_source"] = (args.terms_from if args.terms_from
                                              else "built-in symbolic pattern only")
    after = sum(1 for s_ in t.get("segments", []) if s_.get("notation"))
    emit(t, args.out)
    print(f"mpk: notation-bearing {before} -> {after} ({changed} segment(s) "
          f"changed) using {len(terms)} deck term(s)", file=sys.stderr)


# ── the review file ────────────────────────────────────────────────────────
#
# A human types this while listening, so it has to be typeable: one line,
# three fields, pipe-separated.
#
#     MM:SS | kind | comment
#
# kinds:  fix    heard -> should be
#         check  needs a second look; states what to look at
#         ok     listened, correct as transcribed
#         ask    cannot decide; needs someone else
#         skim   deliberately not listened to — a stated gap, not a verdict
#
# `all` in the time column applies to the whole transcript. Blank lines and
# lines beginning with # are ignored, as are continuation lines (a leading
# pipe), which exist so a long comment can wrap without breaking the format.

_REVIEW_KINDS = {"fix", "check", "ok", "ask", "skim"}
_ARROW = re.compile(r"\s*->\s*")
_BRACKET = re.compile(r"\s*\[[^\]]*\]\s*$")


def _parse_clock(tok: str) -> float | None:
    """'06:28' or '6:28' or '388.1' -> seconds. 'all' -> None."""
    tok = tok.strip()
    if tok.lower() == "all":
        return None
    if ":" in tok:
        parts = tok.split(":")
        try:
            nums = [float(p) for p in parts]
        except ValueError:
            die(f"cannot read a time from {tok!r}")
        if len(nums) == 2:
            return nums[0] * 60 + nums[1]
        if len(nums) == 3:
            return nums[0] * 3600 + nums[1] * 60 + nums[2]
        die(f"cannot read a time from {tok!r}")
    try:
        return float(tok)
    except ValueError:
        die(f"cannot read a time from {tok!r}")


def parse_review(path: str) -> list[dict]:
    """Read a review file into a list of comments. Refuses to guess."""
    out: list[dict] = []
    for lineno, raw in enumerate(Path(path).read_text(encoding="utf-8")
                                 .splitlines(), 1):
        line = raw.strip()
        if not line or line.startswith("#"):
            continue
        if line.startswith("|"):                      # wrapped continuation
            if not out:
                die(f"{path}:{lineno}: continuation line before any comment")
            out[-1]["comment"] += " " + line.lstrip("|").strip()
            continue
        parts = [p.strip() for p in line.split("|")]
        if len(parts) < 3:
            die(f"{path}:{lineno}: expected 'MM:SS | kind | comment', got {raw!r}")
        when, kind, comment = parts[0], parts[1].lower(), "|".join(parts[2:]).strip()
        if kind not in _REVIEW_KINDS:
            die(f"{path}:{lineno}: unknown kind {kind!r}. "
                f"Use one of: {', '.join(sorted(_REVIEW_KINDS))}")
        out.append({"line": lineno, "at": _parse_clock(when), "at_text": when,
                    "kind": kind, "comment": comment})
    if not out:
        die(f"{path} has no review lines")
    return out


def _seg_at(segs: list[dict], t: float) -> dict | None:
    """The segment containing t, or the nearest one if t lands in a pause.

    A reviewer reads the clock off a player, so the time can fall a little
    outside the segment it refers to — and pause comments land in silence by
    definition. Nearest-by-distance is the only sane reading."""
    for s in segs:
        if s.get("start", 0) <= t <= s.get("end", 0):
            return s
    if not segs:
        return None
    return min(segs, key=lambda s: min(abs(s.get("start", 0) - t),
                                       abs(s.get("end", 0) - t)))


def _norm(w: str) -> str:
    return re.sub(r"[^a-z0-9_]", "", w.lower())


def _find_run(words: list[dict], phrase: str) -> tuple[int, int] | None:
    """Locate a consecutive run of words matching `phrase`. Returns [i, j).

    Matching is on letters and digits only, and word-for-word first. If that
    fails it falls back to comparing the phrase with its spaces removed against
    a joined run of tokens — because a reviewer writes what they read on the
    page ("water-independent") while the recogniser may hold it as two tokens
    ("water", "-independent"). Requiring the reviewer to know the tokenisation
    would make the review file unwritable."""
    have = [_norm(w.get("word", "")) for w in words]
    want = [_norm(x) for x in phrase.split() if _norm(x)]
    if not want:
        return None
    for i in range(len(have) - len(want) + 1):
        if have[i:i + len(want)] == want:
            return i, i + len(want)
    glued = "".join(want)                      # tokenisation-blind fallback
    for i in range(len(have)):
        acc = ""
        for j in range(i, min(i + 8, len(have))):
            acc += have[j]
            if acc == glued:
                return i, j + 1
            if len(acc) > len(glued):
                break
    return None


def apply_fix(seg: dict, heard: str, should_be: str) -> dict:
    """Replace a run of words. Timings are preserved, never invented.

    The audio did not change — only our label for it. So the replacement spans
    exactly the first token's start to the last token's end, and no interior
    boundary is created. When several words collapse into one (\"sum of 4\" ->
    \"semaphore\"), the merged word owns the whole span; when one word becomes
    several, they share it evenly, which is a guess and is flagged as one."""
    words = seg.get("words") or []
    run = _find_run(words, heard)
    if run is None:
        return {"ok": False,
                "why": f"cannot find {heard!r} in {seg.get('segment_id', '?')}"}
    i, j = run
    start, end = words[i]["start"], words[j - 1]["end"]
    probs = [w.get("probability", 0) for w in words[i:j]]
    new_tokens = should_be.split()
    # Carry trailing punctuation across. A reviewer writes "sum of 4 ->
    # semaphore" while the token is "4." — dropping the full stop would join
    # two sentences, and sentence boundaries are what Layer 8 reads to find
    # where one idea ends.
    tail = re.search(r"[.,;:?!]+$", words[j - 1].get("word", ""))
    if tail and not re.search(r"[.,;:?!]+$", new_tokens[-1]):
        new_tokens[-1] += tail.group(0)
    if len(new_tokens) == 1:
        replacement = [{"word": new_tokens[0], "start": start, "end": end,
                        "probability": min(probs) if probs else None,
                        "corrected": True}]
        even = False
    else:
        step = (end - start) / len(new_tokens)
        replacement = [{"word": tok,
                        "start": round(start + k * step, 2),
                        "end": round(start + (k + 1) * step, 2),
                        "probability": min(probs) if probs else None,
                        "corrected": True, "timing_estimated": True}
                       for k, tok in enumerate(new_tokens)]
        even = True
    seg["words"] = words[:i] + replacement + words[j:]
    seg["text"] = " ".join(w.get("word", "") for w in seg["words"]).strip()
    return {"ok": True, "start": start, "end": end, "n_before": j - i,
            "n_after": len(new_tokens), "timing_estimated": even}


def cmd_transcript_apply(args):
    """Apply a human review file to a transcript.

    This is the step that turns recorded decisions into the artifact Layer 8
    reads. Everything a human could not settle is carried forward as an open
    question rather than dropped — a review comment that vanishes is worse
    than one that was never written."""
    raw = Path(args.file).read_text(encoding="utf-8")
    try:
        t = json.loads(raw)
    except json.JSONDecodeError as exc:
        die(f"{args.file} is not valid JSON: {exc}")
    t = t.get("payload", t)
    segs = t.get("segments") or []
    if not segs:
        die(f"{args.file} has no segments")

    review = parse_review(args.review)
    who = args.by
    applied, failed, carried, coverage = [], [], [], None

    for c in review:
        if c["at"] is None:                       # the 'all' line
            coverage = c
            continue
        seg = _seg_at(segs, c["at"])
        if seg is None:
            failed.append({**c, "why": "no segment near that time"})
            continue
        sid = seg.get("segment_id", "?")

        if c["kind"] == "fix":
            body = _BRACKET.sub("", c["comment"])
            if not _ARROW.search(body):
                failed.append({**c, "why": "a fix line needs 'heard -> should be'"})
                continue
            heard, should_be = _ARROW.split(body, 1)
            res = apply_fix(seg, heard.strip(), should_be.strip())
            if not res["ok"]:
                # A reviewer reads the clock off a player, and a segment
                # boundary can fall mid-phrase. Look either side before
                # declaring the line unusable.
                k = segs.index(seg)
                for n in (k - 2, k - 1, k + 1, k + 2):
                    if 0 <= n < len(segs):
                        alt = apply_fix(segs[n], heard.strip(), should_be.strip())
                        if alt["ok"]:
                            seg, sid, res = segs[n], segs[n].get("segment_id", "?"), alt
                            res["moved"] = True
                            break
            if not res["ok"]:
                failed.append({**c, "why": res["why"]})
                continue
            seg.setdefault("corrections", []).append({
                "from": heard.strip(), "to": should_be.strip(),
                "span": {"start": res["start"], "end": res["end"],
                         "note": "merged token spans first token start to last "
                                 "token end; no interior boundary invented"
                                 if not res["timing_estimated"] else
                                 "one token expanded to several; interior "
                                 "boundaries are evenly spaced ESTIMATES"},
                "basis": f"human review, {args.review}:{c['line']}",
                "status": "confirmed", "confirmed_by": who,
                "confirmed_how": "listened in the review page"})
            applied.append({**c, "segment_id": sid, **res})

        elif c["kind"] in ("check", "ask"):
            seg.setdefault("flags", []).append({
                "code": f"open_question_{c['kind']}", "severity": "warn",
                "note": f"{c['comment']}  [reviewer {who}, "
                        f"{args.review}:{c['line']}]",
                "needs_eye_check": True, "raised_at_s": c["at"]})
            carried.append({**c, "segment_id": sid})

        elif c["kind"] == "ok":
            seg["reviewed"] = {"by": who, "verdict": "correct as transcribed",
                               "note": c["comment"]}
            seg["flags"] = [f for f in (seg.get("flags") or [])
                            if f.get("code") != "low_confidence_words"]

        elif c["kind"] == "skim":
            seg.setdefault("flags", []).append({
                "code": "not_reviewed", "severity": "info",
                "note": f"{c['comment']}  [reviewer {who}]",
                "needs_eye_check": True})

    md = t.setdefault("metadata", {})
    md["extraction_path"] = md.get("extraction_path", "") + "+review_applied"
    md["review_file"] = args.review
    md["reviewed_by"] = who
    md["corrections_applied"] = len(applied)
    md["open_questions"] = len(carried)
    md["review_coverage"] = (
        {"kind": coverage["kind"], "note": coverage["comment"]} if coverage
        else {"kind": "unstated",
              "note": "the review file has no 'all' line, so coverage of the "
                      "unlisted segments is unknown"})
    md["note"] = ("Human-reviewed transcript. Corrections are status "
                  "'confirmed' — a person listened. Timings were never altered: "
                  "a replacement spans exactly what it replaced. Open questions "
                  "from the review are carried as flags, not resolved here.")

    emit(t, args.out)
    print(f"mpk: {len(applied)} correction(s) applied, {len(carried)} open "
          f"question(s) carried, {len(failed)} line(s) FAILED", file=sys.stderr)
    for a in applied:
        print(f"  ok    {a['at_text']:>6}  {a['segment_id']}  "
              f"{a['n_before']}->{a['n_after']} word(s)"
              + ("  [timing estimated]" if a["timing_estimated"] else ""),
              file=sys.stderr)
    for c in carried:
        print(f"  open  {c['at_text']:>6}  {c['segment_id']}  {c['kind']}",
              file=sys.stderr)
    for f in failed:
        print(f"  FAIL  {f['at_text']:>6}  line {f['line']}: {f['why']}",
              file=sys.stderr)
    if failed and not args.keep_going:
        die(f"{len(failed)} review line(s) could not be applied. Fix them, or "
            f"pass --keep-going to write the transcript without them.")


def cmd_transcript_check(args):
    raw = Path(args.file).read_text(encoding="utf-8")
    try:
        t = json.loads(raw).get("payload", json.loads(raw))
    except json.JSONDecodeError as exc:
        die(f"{args.file} is not valid JSON: {exc}")
    segs = t.get("segments", [])
    problems, notes = [], []
    if not segs:
        problems.append("no segments")
    n_words = sum(len(s.get("words") or []) for s in segs)
    if n_words == 0:
        problems.append("no word-level timings anywhere — VGR-05 cannot be met "
                        "and Layer 8 has nothing to bind to")
    for s in segs:
        sid = s.get("segment_id", "?")
        if not (s.get("words") or []):
            notes.append(f"{sid}: no word timings")
        for w in (s.get("words") or []):
            if w.get("start") is None or w.get("end") is None:
                problems.append(f"{sid}: word {w.get('word')!r} missing a timing")
            elif w["end"] < w["start"]:
                problems.append(f"{sid}: word {w.get('word')!r} ends before it starts")
    for a, b in zip(segs, segs[1:]):
        if b.get("start", 0) < a.get("end", 0):
            notes.append(f"{a.get('segment_id')} overlaps {b.get('segment_id')}")
    notation = [s for s in segs if s.get("notation")]
    print(f"segments: {len(segs)}  words: {n_words}  "
          f"notation-bearing: {len(notation)}")
    for p_ in problems:
        print(f"  FAIL  {p_}")
    for n_ in notes[:20]:
        print(f"  note  {n_}")
    if len(notes) > 20:
        print(f"  note  … and {len(notes)-20} more")
    print(f"\n{len(problems)} failure(s), {len(notes)} note(s)")
    sys.exit(1 if problems else 0)


# ════════════════════════════════════════════════════════════ audio / video
def _ff(cmd: list[str]) -> str:
    return subprocess.run(cmd, capture_output=True, text=True).stdout


def cmd_audio_extract(args):
    if not have("ffmpeg"):
        die("ffmpeg not found")
    # RC-003: the delivery master is 48 kHz stereo and must never be the
    # 16 kHz copy that speech recognition wants.
    subprocess.run(["ffmpeg", "-y", "-i", args.file, "-vn",
                    "-ar", "48000", "-ac", "2", "-c:a", "pcm_s16le", args.out],
                   check=True, stderr=subprocess.DEVNULL)
    print(f"wrote {args.out} — 48 kHz stereo master (RC-003 delivery path)",
          file=sys.stderr)


def cmd_audio_asr(args):
    if not have("ffmpeg"):
        die("ffmpeg not found")
    subprocess.run(["ffmpeg", "-y", "-i", args.file, "-vn",
                    "-ar", "16000", "-ac", "1", "-c:a", "pcm_s16le", args.out],
                   check=True, stderr=subprocess.DEVNULL)
    print(f"wrote {args.out} — 16 kHz mono ASR copy. Never ship this as the "
          f"master (RC-003).", file=sys.stderr)


def cmd_audio_probe(args):
    if not have("ffmpeg"):
        die("ffmpeg not found")
    j = json.loads(_ff(["ffprobe", "-v", "quiet", "-print_format", "json",
                        "-show_streams", "-show_format", args.file]) or "{}")
    a = next((s for s in j.get("streams", []) if s.get("codec_type") == "audio"), {})
    vd = subprocess.run(["ffmpeg", "-i", args.file, "-af", "volumedetect",
                         "-f", "null", "-"], capture_output=True, text=True).stderr
    mean = re.search(r"mean_volume:\s*(-?[\d.]+) dB", vd)
    peak = re.search(r"max_volume:\s*(-?[\d.]+) dB", vd)
    emit({"file": Path(args.file).name,
          "codec": a.get("codec_name"), "sample_rate_hz": a.get("sample_rate"),
          "channels": a.get("channels"), "bitrate_bps": a.get("bit_rate"),
          "duration_s": j.get("format", {}).get("duration"),
          "mean_volume_db": float(mean.group(1)) if mean else None,
          "peak_volume_db": float(peak.group(1)) if peak else None,
          "note": "Layer 7 input. Cite TGT-005…010 rather than repeating these numbers."},
         args.out)


def cmd_video_probe(args):
    if not have("ffprobe"):
        die("ffprobe not found")
    j = json.loads(_ff(["ffprobe", "-v", "quiet", "-print_format", "json",
                        "-show_streams", "-show_format", args.file]) or "{}")
    v = next((s for s in j.get("streams", []) if s.get("codec_type") == "video"), {})
    a = next((s for s in j.get("streams", []) if s.get("codec_type") == "audio"), {})
    fr = v.get("avg_frame_rate", "0/1")
    try:
        n, d = fr.split("/"); fps = round(int(n) / int(d), 3) if int(d) else None
    except Exception:
        fps = None
    emit({"file": Path(args.file).name,
          "resolution": f"{v.get('width')}x{v.get('height')}",
          "declared_fps": fps, "video_codec": v.get("codec_name"),
          "profile": v.get("profile"),
          "video_bitrate_mbps": (round(int(v["bit_rate"]) / 1e6, 2)
                                 if v.get("bit_rate") else None),
          "audio_codec": a.get("codec_name"), "audio_channels": a.get("channels"),
          "audio_sample_rate_hz": a.get("sample_rate"),
          "duration_s": j.get("format", {}).get("duration"),
          "note": "declared_fps is what the container reports. It does NOT prove the "
                  "picture changes that often — see 'mpk video uniquefps' (TGT-013)."},
         args.out)


def _render_deck(path: str, dpi: int = 100) -> list:
    """Render a deck to PNGs in a temp dir. Reuses the same soffice -> pdftoppm
    path as `mpk deck render`, at a lower dpi because these are only ever
    compared at 96x54."""
    soffice = next((c for c in ("soffice", "libreoffice") if have(c)), None)
    if not soffice:
        die("LibreOffice not found — needed to render the deck for slide matching. "
            "Install it, or run without --deck.")
    if not have("pdftoppm"):
        die("pdftoppm not found (poppler-utils) — needed to turn the rendered PDF "
            "into images. Install it, or run without --deck.")
    tmp = Path(tempfile.mkdtemp(prefix="mpk-render-"))
    stem = Path(path).stem
    subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir",
                    str(tmp), path], check=True, stdout=subprocess.DEVNULL,
                   stderr=subprocess.DEVNULL)
    pdf = tmp / (stem + ".pdf")
    if not pdf.exists():
        die(f"LibreOffice produced no PDF for {path}")
    subprocess.run(["pdftoppm", "-png", "-r", str(dpi), str(pdf), str(tmp / "s")],
                   check=True)
    pages = sorted(tmp.glob("s-*.png"))
    if not pages:
        die(f"no pages rendered from {path}")
    return pages


def _frames(path: str, fps: float, w: int, h: int, rgb: bool = False):
    """Every frame at `fps`, w x h, as one numpy array.

    Raw video straight out of ffmpeg — no image library for the greyscale path.
    numpy arrives with faster-whisper, so this adds nothing the project did not
    already depend on."""
    try:
        import numpy as np
    except ImportError:
        die("numpy is required for this command (it ships with faster-whisper): "
            "pip install numpy")
    if not have("ffmpeg"):
        die("ffmpeg is required — install it with your package manager "
            "(apt install ffmpeg)")
    fmt, chan = ("rgb24", 3) if rgb else ("gray", 1)
    p = subprocess.run(
        ["ffmpeg", "-v", "error", "-i", path, "-vf", f"fps={fps},scale={w}:{h}",
         "-pix_fmt", fmt, "-f", "rawvideo", "-"], capture_output=True)
    if p.returncode != 0:
        die(f"ffmpeg failed reading {path}: {p.stderr.decode()[:400]}")
    n = len(p.stdout) // (w * h * chan)
    if n < 2:
        die(f"only {n} frame(s) decoded — is {path} a video?")
    a = np.frombuffer(p.stdout[:n * w * h * chan], dtype="uint8")
    return (a.reshape(n, h, w, 3) if rgb else a.reshape(n, h * w)), np


def _thumb(path: str, t: float, width: int) -> str:
    """One frame at t, as a data: URI. Empty string if it cannot be read."""
    p = subprocess.run(
        ["ffmpeg", "-v", "error", "-ss", f"{max(t, 0):.3f}", "-i", path,
         "-frames:v", "1", "-vf", f"scale={width}:-2", "-q:v", "6",
         "-f", "image2", "-vcodec", "mjpeg", "-"], capture_output=True)
    if p.returncode != 0 or not p.stdout:
        return ""
    return "data:image/jpeg;base64," + base64.b64encode(p.stdout).decode()


def _png_uri(path, width: int = 480) -> str:
    try:
        from PIL import Image
    except ImportError:
        return ""
    import io
    im = Image.open(path).convert("RGB")
    im = im.resize((width, max(1, round(im.height * width / im.width))))
    b = io.BytesIO(); im.save(b, "JPEG", quality=72)
    return "data:image/jpeg;base64," + base64.b64encode(b.getvalue()).decode()


def _norm(x, np):
    x = x - x.mean()
    sd = x.std()
    return x / sd if sd > 1e-6 else x


def _build_review(data: dict, out_path: str, template: str):
    """Emit the review page from the data we just produced.

    The page is a complete HTML file whose embedded JavaScript renders whatever
    is in its data block — so the table and the JSON are the same object and
    cannot disagree. This helper only substitutes the block."""
    tpl = resolve_template(template).read_text(encoding="utf-8")
    html, n = _DATA_BLOCK.subn(
        lambda m: m.group(1) + json.dumps(data, indent=1, ensure_ascii=False)
        + m.group(2), tpl)
    if n != 1:
        die(f"expected exactly one data block in {template}")
    Path(out_path).write_text(html, encoding="utf-8")
    print(f"wrote {out_path}  ({Path(out_path).stat().st_size:,} bytes) — {template}",
          file=sys.stderr)


def cmd_video_slidechanges(args):
    """What is on screen, second by second — measured from the picture.

    Three things this answers that nothing else can:

    1. WHEN the picture changes. A pause is NOT a slide change: on V017 the
       professor stops speaking at 4:09 and the slide turns at 4:13.75. He
       finishes the thought, pauses, then advances. Deriving slide windows from
       transcript pauses puts every boundary about four seconds early, and the
       symptom surfaces two layers later looking like a sync bug.

    2. WHICH slide it is. The deck is rendered and each window's middle frame is
       matched against it. Without this the mapping is an assumption — and on
       V017 that assumption was WRONG: deck slides 1 and 2 never appear in the
       video at all, and the first 82 seconds is a presenter composition that
       exists in no deck.

    3. WHERE HE POINTED. The source video carries a hand-placed orange arrow
       marking what he is talking about. It is extracted as ground truth, so a
       later layer's focus map can be CHECKED rather than argued about.

    Two metrics run together, because one is blind where the other is not. Mean
    absolute difference catches a slide turn, which moves most of the picture.
    It cannot see a logo cross-dissolving on a white card — that is 1% of pixels
    and scores 2 against a threshold of 15. Percentage-of-pixels-changed catches
    exactly that case. V017's whole intro sequence is invisible to the first
    metric and obvious to the second."""
    frames, np = _frames(args.file, args.fps, args.width, args.height)
    f32 = frames.astype("float32")
    d = np.abs(np.diff(f32, axis=0))
    mean_d = d.mean(1)
    frac_d = (d > args.pixel_delta).mean(1) * 100.0
    t_of = lambda i: (i + 1) / args.fps
    dur = len(frames) / args.fps

    # ---- events: either metric is enough. A cross-dissolve never trips the
    #      first; a slide turn always trips both.
    hit = [i for i in range(len(mean_d))
           if mean_d[i] >= args.weak_threshold or frac_d[i] >= args.frac_threshold]
    clusters, cur = [], []
    for i in hit:
        if cur and t_of(i) - t_of(cur[-1]) > args.min_gap:
            clusters.append(cur); cur = []
        cur.append(i)
    if cur:
        clusters.append(cur)

    changes = []
    for c in clusters:
        peak = max(c, key=lambda i: float(mean_d[i]))
        m, f = float(mean_d[peak]), float(max(frac_d[i] for i in c))
        changes.append({
            "t": round(t_of(peak), 3),
            "score": round(m, 2),
            "pct_pixels_changed": round(f, 2),
            "seen_by": ("both" if m >= args.weak_threshold and f >= args.frac_threshold
                        else "mean" if m >= args.weak_threshold else "pixel_fraction"),
            "tier": "strong" if m >= args.threshold else "weak",
            "spans_s": [round(t_of(c[0]), 3), round(t_of(c[-1]), 3)],
            "frames_in_cluster": len(c),
            "confirmed": None, "verdict": None})

    # ---- windows: what is held between the changes
    bounds = [0.0] + [ch["spans_s"][1] for ch in changes]
    if bounds[-1] < dur - 0.5:
        bounds.append(round(dur, 3))
    windows = []
    for a, b in zip(bounds, bounds[1:]):
        if b - a < args.min_window:
            continue
        windows.append({"start": round(a, 3), "end": round(b, 3),
                        "duration_s": round(b - a, 3),
                        "slide_id": None, "match_score": None,
                        "runner_up": None, "runner_up_score": None,
                        "kind": "unknown", "confirmed": None})

    # ---- which slide, if a deck was given
    deck_pages = []
    if args.deck:
        pages = _render_deck(args.deck, args.render_dpi)
        deck_pages = [str(p) for p in pages]
        try:
            from PIL import Image
        except ImportError:
            die("Pillow is required to match rendered slides: pip install pillow")
        W2, H2 = args.match_width, args.match_height
        refs = [_norm(np.asarray(Image.open(p).convert("L")
                                 .resize((W2, H2), Image.BILINEAR),
                                 dtype="float32").ravel(), np) for p in pages]
        mid, _ = _frames(args.file, args.fps, W2, H2)
        midf = mid.astype("float32")
        for w in windows:
            i = min(int(((w["start"] + w["end"]) / 2) * args.fps), len(midf) - 1)
            fv = _norm(midf[i], np)
            sc = [float((fv * r).mean()) for r in refs]
            order = sorted(range(len(sc)), key=lambda k: sc[k], reverse=True)
            best = order[0]
            w["match_score"] = round(sc[best], 3)
            w["runner_up"] = f"{Path(args.deck).stem}#{order[1] + 1}" if len(sc) > 1 else None
            w["runner_up_score"] = round(sc[order[1]], 3) if len(sc) > 1 else None
            if sc[best] >= args.match_min:
                w["slide_id"] = f"{args.deck_id or Path(args.deck).stem}_s{best + 1:02d}"
                w["slide_number"] = best + 1
                w["kind"] = "slide"
            else:
                w["kind"] = "not_in_deck"
                w["note"] = ("No deck slide matches this window. It is a title card, an "
                             "outro, or a composition built for the video that exists in "
                             "no deck. NOT an error — on V017 the first 82 seconds are "
                             "exactly this.")

        # ---- merge neighbours showing the SAME picture.
        # An event is not always a new slide. A focus arrow moving, or the
        # presenter shifting, changes enough pixels to trip the detector while
        # the slide underneath never turns. Two adjacent windows resolving to
        # one slide_id ARE one window, and saying otherwise hands Layer 8 four
        # boundaries where the deck has one.
        merged, dropped = [], 0
        for w in windows:
            prev = merged[-1] if merged else None
            same = False
            if prev is not None:
                if w["slide_id"] and prev["slide_id"] == w["slide_id"]:
                    same = True
                elif w["kind"] == "not_in_deck" == prev["kind"]:
                    ia = min(int(((prev["start"] + prev["end"]) / 2) * args.fps),
                             len(midf) - 1)
                    ib = min(int(((w["start"] + w["end"]) / 2) * args.fps), len(midf) - 1)
                    same = float((_norm(midf[ia], np) * _norm(midf[ib], np)).mean()) > \
                           args.merge_min
            if same:
                prev["end"] = w["end"]
                prev["duration_s"] = round(prev["end"] - prev["start"], 3)
                prev["match_score"] = max(prev["match_score"] or -1, w["match_score"] or -1)
                prev["merged_from"] = prev.get("merged_from", 1) + 1
                dropped += 1
            else:
                merged.append(w)
        windows = merged
        if dropped:
            print(f"mpk: merged {dropped} window(s) that showed the same picture "
                  f"as their neighbour", file=sys.stderr)

    # ---- where he pointed
    focus = []
    if args.arrow:
        # Only hunt on windows that matched a deck slide. The arrow lives on the
        # slides; everywhere else, "orange" is a logo. On V017 the KLE mark (red)
        # and the ekLakshya dots (orange) produced three false runs in the intro,
        # which then read as the focus map disagreeing with ground truth when the
        # ground truth was the thing that was wrong.
        deck_spans = [(w["start"], w["end"]) for w in windows if w.get("slide_id")]
        rgbf, _ = _frames(args.file, args.arrow_fps, args.width, args.height, rgb=True)
        a = rgbf.astype(int)
        R, G, B = a[..., 0], a[..., 1], a[..., 2]
        orange = ((R > args.arrow_r) & (G > 70) & (G < 175) & (B < 90)
                  & ((R - B) > 110))
        count = orange.reshape(len(a), -1).sum(1)
        present = count >= args.arrow_min_px
        if deck_spans:
            for k in range(len(present)):
                t = k / args.arrow_fps
                if not any(a <= t <= b for a, b in deck_spans):
                    present[k] = False
        i = 0
        while i < len(present):
            if present[i]:
                j = i
                while j < len(present) and present[j]:
                    j += 1
                if (j - i) / args.arrow_fps >= args.arrow_min_s:
                    xs, ys = [], []
                    for k in range(i, j):
                        yy, xx = np.where(orange[k])
                        xs.append(xx.mean() / args.width); ys.append(yy.mean() / args.height)
                    focus.append({
                        "start": round(i / args.arrow_fps, 2),
                        "end": round(j / args.arrow_fps, 2),
                        "duration_s": round((j - i) / args.arrow_fps, 2),
                        "x": round(float(np.mean(xs)), 3),
                        "y_from": round(float(np.min(ys)), 3),
                        "y_to": round(float(np.max(ys)), 3),
                        "note": "x says which column, y says roughly which line"})
                i = j
            else:
                i += 1

    strong = [c for c in changes if c["tier"] == "strong"]
    out = {
      "video_id": args.video_id or Path(args.file).stem,
      "metadata": {
        "source_file": Path(args.file).name,
        "deck_file": Path(args.deck).name if args.deck else None,
        "duration_s": round(dur, 3),
        "sample_fps": args.fps, "sample_size": [args.width, args.height],
        "metrics": {
          "mean_abs_diff": "mean absolute difference between consecutive greyscale "
                           "frames, 0-255. Catches a slide turn.",
          "pct_pixels_changed": f"percentage of pixels changing by more than "
                                f"{args.pixel_delta} levels. Catches a cross-dissolve or "
                                f"a small logo swap, which the mean cannot see."},
        "thresholds": {"strong": args.threshold, "weak": args.weak_threshold,
                       "pct_pixels": args.frac_threshold, "min_gap_s": args.min_gap,
                       "slide_match_min": args.match_min},
        "noise_floor": {"mean_p50": round(float(np.percentile(mean_d, 50)), 3),
                        "mean_p99": round(float(np.percentile(mean_d, 99)), 3),
                        "pct_p99": round(float(np.percentile(frac_d, 99)), 3)},
        "counts": {"changes": len(changes), "strong": len(strong),
                   "weak": len(changes) - len(strong), "windows": len(windows),
                   "windows_matched": sum(1 for w in windows if w["slide_id"]),
                   "windows_not_in_deck": sum(1 for w in windows
                                              if w["kind"] == "not_in_deck"),
                   "focus_runs": len(focus)},
        "deck_pages_rendered": len(deck_pages),
        "tier_meaning": "How much of the picture moved, NOTHING MORE. 'weak' does not "
                        "mean 'not a slide change': on V017 the 01:22.75 turn scored "
                        "17.00 and was confirmed by eye, while the highest-scoring event "
                        "in the whole video (160.15) was the opening fade, which is not a "
                        "slide change at all.",
        "threshold_provenance": "strong=15.0, set below the lowest CONFIRMED slide change "
                                "measured so far (17.00, V017 01:22.75). One video is thin "
                                "evidence for a constant — revisit after V018 (OBS-035).",
        "extraction_path": "mpk_video_slidechanges", "mpk_version": __version__,
        "note": "WHERE the picture changes, WHICH slide it is, and WHERE he pointed. "
                "Never WHY. Every window still needs a human verdict in the review page."},
      "changes": changes, "windows": windows, "focus_ground_truth": focus}

    if args.thumbs:
        for ch in changes:
            ch["before"] = _thumb(args.file, ch["spans_s"][0] - args.thumb_lead,
                                  args.thumb_width)
            ch["after"] = _thumb(args.file, ch["spans_s"][1] + args.thumb_lead,
                                 args.thumb_width)
        for w in windows:
            w["frame"] = _thumb(args.file, (w["start"] + w["end"]) / 2, args.thumb_width)
            n = w.get("slide_number")
            w["render"] = _png_uri(deck_pages[n - 1], args.thumb_width) if n else ""

    emit(out, args.out)
    mm = lambda v: f"{int(v // 60):02d}:{v % 60:05.2f}"
    print(f"mpk: {len(changes)} change(s), {len(windows)} window(s), "
          f"{len(focus)} focus run(s)", file=sys.stderr)
    for w in windows:
        tag = w["slide_id"] or ("NOT IN DECK" if w["kind"] == "not_in_deck" else "—")
        sc = f"{w['match_score']:.3f}" if w["match_score"] is not None else "    -"
        print(f"  {mm(w['start'])}–{mm(w['end'])}  {w['duration_s']:7.2f}s  "
              f"{tag:16} match {sc}", file=sys.stderr)
    if args.html:
        _build_review(out, args.html, "slidechange-review")


def cmd_video_uniquefps(args):
    """TGT-013. A 30 fps file made of 8 unique frames per second reports 30 to
    ffprobe and passes TGT-002, while still looking choppy — because encoding
    cannot add motion that was never rendered (RC-001)."""
    if not have("ffmpeg"):
        die("ffmpeg not found")
    dur = _ff(["ffprobe", "-v", "quiet", "-show_entries", "format=duration",
               "-of", "csv=p=0", args.file]).strip()
    err = subprocess.run(["ffmpeg", "-i", args.file, "-vf",
                          f"mpdecimate=hi={args.hi}:lo={args.lo}",
                          "-loglevel", "info", "-f", "null", "-"],
                         capture_output=True, text=True).stderr
    m = re.search(r"frame=\s*(\d+)", err[::-1][:400][::-1]) or \
        re.findall(r"frame=\s*(\d+)", err)
    kept = int(m[-1]) if isinstance(m, list) and m else (int(m.group(1)) if m else None)
    d = float(dur) if dur else None
    eff = round(kept / d, 2) if kept and d else None
    emit({"file": Path(args.file).name, "duration_s": d,
          "unique_frames": kept, "effective_unique_fps": eff,
          "threshold": args.threshold,
          "verdict": (None if eff is None else
                      "PASS" if eff >= args.threshold else "FAIL"),
          "enforcement": "advisory",
          "note": "TGT-013 — measured over the whole file here. The target is "
                  "measured over ACTIVE BEAT WINDOWS, so deliberately still "
                  "content is not penalised; a whole-file figure is a floor, "
                  "not the gate."},
         args.out)


# ════════════════════════════════════════════════════════════ check manifest

# The vocabulary Layer 5 resolves into (DEC-002). Fixed on purpose: without it
# a prompt invents a name per context, one video carries two names for one idea,
# Layer 6 makes two bindings, and the same event appears in two colours. The list
# constrains MEANING, never FORM — shape comes from Layer 3's `type`, so an arrow
# stays an arrow and a code line stays a code line.
L5_VOCAB = {
    "process_actor", "resource", "counter", "capacity", "lock", "state_flag",
    "enqueue", "dequeue", "request", "assignment", "access",
    "wait_point", "signal_point", "guard",
    "title", "body_text", "code_block", "code_line", "matrix", "matrix_cell",
    "label", "annotation", "panel", "list_item",
    "background", "footer_band", "logo", "slide_number",
}
L5_ROLES = {"title_card", "recap", "motivation", "walkthrough",
            "worked_example", "summary", "exercise", "transition"}

# Chrome is settled by GEOMETRY, not by narration. Nobody says "and here is the
# footer band" — it is a footer band because of where it sits and what it
# repeats, which Layer 3 already established in chrome_pattern. Demanding a
# narration citation for these would force a prompt either to leave real
# classifications unresolved or to invent a quote, and inventing a quote is far
# worse than the rule was ever worth.
L5_NO_CITATION_NEEDED = {"background", "footer_band", "logo", "slide_number", "title"}


def cmd_check_representation(args):
    """Validate a Layer 5 artifact against the layers it came from.

    VGR-07 says nothing narrated may be silently dropped from the plan. Until
    now that was a rule with no enforcement — a representation could omit half
    the deck and still look complete. This is the enforcement."""
    def load(p, what):
        try:
            d = json.loads(Path(p).read_text(encoding="utf-8"))
        except FileNotFoundError:
            die(f"{what} not found: {p}")
        except json.JSONDecodeError as exc:
            die(f"{what} at {p} is not valid JSON: {exc}")
        return d.get("payload", d)

    r = load(args.file, "representation")
    fails, notes = [], []

    elements = r.get("elements") or []
    fmap = sorted(r.get("focus_map") or [], key=lambda w: w.get("start", 0))
    if not elements:
        fails.append("no elements")
    if not fmap:
        fails.append("no focus_map — Layer 8 has nothing to bind against")

    # ---- 1 · element coverage against the manifest (VGR-07)
    if args.manifest:
        man = load(args.manifest, "manifest")
        want = {a["element_id"] for s in man.get("slides", []) for a in s.get("assets", [])}
        got = [x.get("element_id") for x in elements]
        seen = set(got)
        missing = sorted(want - seen)
        extra = sorted(seen - want)
        dupes = sorted({i for i in got if got.count(i) > 1})
        if missing:
            fails.append(f"{len(missing)} manifest element(s) never appear: "
                         f"{', '.join(missing[:8])}{' …' if len(missing) > 8 else ''}")
        if extra:
            fails.append(f"{len(extra)} element(s) invented — not in the manifest: "
                         f"{', '.join(extra[:8])}")
        if dupes:
            fails.append(f"{len(dupes)} element(s) appear more than once: "
                         f"{', '.join(dupes[:8])}")

    # ---- 2 · every resolution cites its evidence
    for x in elements:
        if x.get("status") == "resolved":
            b = x.get("basis") or {}
            cited = b.get("quote") and b.get("segment") and b.get("at") is not None
            if not cited and x.get("semantic_type") not in L5_NO_CITATION_NEEDED:
                fails.append(f"{x.get('element_id')}: resolved with no citation "
                             f"(needs quote, segment and timestamp)")
            st = x.get("semantic_type")
            if st and st not in L5_VOCAB:
                added = {v.get("term") for v in (r.get("vocabulary_additions") or [])}
                if st in added:
                    notes.append(f"{x.get('element_id')}: '{st}' is a declared addition")
                else:
                    fails.append(f"{x.get('element_id')}: '{st}' is outside the fixed "
                                 f"vocabulary and is not declared in "
                                 f"vocabulary_additions (DEC-002)")
        elif x.get("status") not in ("no_narration", "unresolved"):
            fails.append(f"{x.get('element_id')}: status "
                         f"{x.get('status')!r} is not one of "
                         f"resolved / no_narration / unresolved")

    # ---- 3 · the focus map must tile the video
    for a, b in zip(fmap, fmap[1:]):
        gap = round(b.get("start", 0) - a.get("end", 0), 3)
        if abs(gap) > args.tolerance:
            fails.append(f"focus map {'gap' if gap > 0 else 'overlap'} of "
                         f"{abs(gap):.2f}s at {a.get('end')}s")
    for w in fmap:
        if w.get("role") not in L5_ROLES:
            fails.append(f"focus window at {w.get('start')}s has role "
                         f"{w.get('role')!r}, outside the fixed list")
        if not w.get("basis"):
            notes.append(f"focus window at {w.get('start')}s has no basis")

    # ---- 4 · nothing here may decide appearance, position or timing
    # Whole words only. Substring matching flagged "Spans changes windows" for
    # `span` and "releasing the mutex" for `easing` — two false alarms in the
    # first real run, which is how a check trains people to ignore it.
    # "span" as a bare verb is ordinary English — the first real run tripped on
    # "no pointer run covers this span". Only the layout sense matters, and that
    # always arrives attached to a grid word.
    BANNED = ("colour", "color", "rgb", "px", "grid", "column",
              "col_span", "row_span", "duration_s", "easing", "hex")
    blob = json.dumps({"e": elements, "f": fmap}).lower()
    hits = [b for b in BANNED if re.search(rf"\b{b}\b", blob)]
    if hits:
        notes.append("possible appearance/layout values present — Layer 5 decides "
                     f"meaning only. Look for: {', '.join(hits)}")

    # ---- 5 · does a focus window name a slide that is on screen then?
    if args.windows:
        wins = load(args.windows, "windows")
        spans = [(w["start"], w["end"], w.get("slide_id"))
                 for w in wins.get("windows", [])]
        for w in fmap:
            sid = w.get("slide")
            if not sid:
                continue
            mid = (w.get("start", 0) + w.get("end", 0)) / 2
            on = next((s for a, b, s in spans if a <= mid <= b), None)
            if on and on != sid:
                fails.append(f"focus window at {w.get('start')}s names {sid} but "
                             f"{on} is on screen then")

        # ---- 6 · does the focus map agree with the hand-placed arrow?
        truth = wins.get("focus_ground_truth") or []
        if truth:
            agree = disagree = 0
            for t in truth:
                mid = (t["start"] + t["end"]) / 2
                w = next((w for w in fmap if w.get("start", 0) <= mid <= w.get("end", 0)),
                         None)
                if w is None:
                    disagree += 1
                elif w.get("subject"):
                    agree += 1
                else:
                    disagree += 1
            notes.append(f"focus arrow: {agree} of {len(truth)} runs land in a window "
                         f"that names a subject, {disagree} do not. The arrow covers "
                         f"only part of the run time — it can confirm a wrong answer "
                         f"but never produce a right one, so disagreement is reported, "
                         f"never corrected")

    print(f"elements: {len(elements)}  focus windows: {len(fmap)}  "
          f"entities: {len(r.get('entities') or [])}", file=sys.stderr)
    for n in notes:
        print(f"  note  {n}", file=sys.stderr)
    for f in fails:
        print(f"  FAIL  {f}", file=sys.stderr)
    print(f"\n{len(fails)} failure(s), {len(notes)} note(s)", file=sys.stderr)
    if fails:
        sys.exit(1)


def cmd_check_manifest(args):
    try:
        m = json.loads(Path(args.file).read_text())
    except json.JSONDecodeError as exc:
        die(f"{args.file} is not valid JSON: {exc}")
    m = m.get("payload", m)
    slides = m.get("slides") or [m]
    problems, notes = [], []

    for s in slides:
        sid = s.get("slide_id", "?")
        for a in s.get("assets", []):
            eid = a.get("element_id", "?")
            if not eid or not re.match(r".+_\d+$", str(eid)):
                problems.append(f"{sid}: element_id {eid!r} does not look "
                                f"OOXML-derived (<slide>_<shape_id>)")
            if a.get("tag") not in ("STATIC", "DYNAMIC"):
                problems.append(f"{eid}: tag {a.get('tag')!r} is not STATIC/DYNAMIC")
            if a.get("semantic_type") is None and not a.get("semantic_type_prior"):
                notes.append(f"{eid}: semantic_type null with no prior recorded")
            if a.get("type") == "connector":
                if not (a.get("properties") or {}).get("direction_verified"):
                    notes.append(f"{eid}: connector direction unverified")
    d = m.get("deck") or {}
    for k in ("entity_inventory", "chrome_pattern", "colour_convention"):
        if k not in d:
            notes.append(f"deck.{k} missing — Layer 6 reads entity_inventory")

    print(f"slides: {len(slides)}  assets: "
          f"{sum(len(s.get('assets', [])) for s in slides)}")
    for p in problems:
        print(f"  FAIL  {p}")
    for n in notes:
        print(f"  note  {n}")
    print(f"\n{len(problems)} failure(s), {len(notes)} note(s)")
    sys.exit(1 if problems else 0)


# ════════════════════════════════════════════════════════════ CLI
def build_parser():
    ap = argparse.ArgumentParser(
        prog="mpk", description="mpk — Media Pipeline Kit. Deterministic utilities "
                                "for the Lecture Alive AI workflow.",
        epilog="Try: mpk deck --help  ·  mpk deck extract --help",
        formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--version", action="version", version=f"mpk {__version__}")
    groups = ap.add_subparsers(dest="group", metavar="<group>")

    # ---- deck
    g = groups.add_parser("deck", help="PowerPoint decks: inspect, extract, reshape")
    gs = g.add_subparsers(dest="cmd", metavar="<command>")

    c = gs.add_parser("info", help="one-line summary per slide")
    c.add_argument("file"); c.set_defaults(fn=cmd_deck_info)

    c = gs.add_parser("extract", help="OOXML shape tree -> raw JSON (Layer 3 step 1)")
    c.add_argument("file")
    c.add_argument("--out", "-o", help="write JSON here (default stdout)")
    c.add_argument("--deck-id", help="id prefix for element ids (default: filename stem)")
    c.add_argument("--slide", type=int, help="extract one slide only")
    c.add_argument("--expect", type=int, help="warn if the slide count differs")
    c.set_defaults(fn=cmd_deck_extract)

    c = gs.add_parser("normalize", help="rescale a deck to another canvas size")
    c.add_argument("file")
    c.add_argument("--like", help="copy the canvas size from this deck")
    c.add_argument("--width", type=int); c.add_argument("--height", type=int)
    c.add_argument("--out", "-o", required=True)
    c.set_defaults(fn=cmd_deck_normalize)

    c = gs.add_parser("merge", help="splice slides from one deck into another")
    c.add_argument("file", help="deck whose slides are inserted")
    c.add_argument("--into", required=True, help="base deck")
    c.add_argument("--at", type=int, help="1-based position (default: append)")
    c.add_argument("--layout", type=int, help="layout index for inserted slides")
    c.add_argument("--out", "-o", required=True)
    c.set_defaults(fn=cmd_deck_merge)

    c = gs.add_parser("render", help="slides -> PNG via LibreOffice")
    c.add_argument("file")
    c.add_argument("--outdir", "-d", default="render")
    c.add_argument("--dpi", type=int, default=150)
    c.set_defaults(fn=cmd_deck_render)

    # ---- review
    g = groups.add_parser("review", help="review pages built from templates")
    gs = g.add_subparsers(dest="cmd", metavar="<command>")
    c = gs.add_parser("build", help="manifest JSON -> self-contained review HTML")
    c.add_argument("file"); c.add_argument("--out", "-o", required=True)
    c.add_argument("--template", "-t", default=None,
                   help=f"template name or path (default: {DEFAULT_TEMPLATE})")
    c.set_defaults(fn=cmd_review_build)
    c = gs.add_parser("templates", help="list available review templates")
    c.set_defaults(fn=cmd_review_templates)

    # ---- transcript
    g = groups.add_parser("transcript",
                          help="transcribe, align, verify and export the timeline")
    gs = g.add_subparsers(dest="cmd", metavar="<command>")

    c = gs.add_parser("build", help="audio -> raw ASR JSON with word timings")
    c.add_argument("file", help="audio or video file (16 kHz mono is enough)")
    c.add_argument("--out", "-o", help="write JSON here (default stdout)")
    c.add_argument("--video-id", help="id for the transcript (default: filename stem)")
    c.add_argument("--model", default="small",
                   help="tiny | base | small | medium | large-v3 (default: small)")
    c.add_argument("--device", default="cpu")
    c.add_argument("--compute", default="int8", help="int8 | float16 | float32")
    c.add_argument("--language", default="en")
    c.add_argument("--beam", type=int, default=5)
    c.add_argument("--pause", type=float, default=1.0,
                   help="report gaps at or above this many seconds (default 1.0)")
    c.add_argument("--low-conf", type=float, default=0.55,
                   help="flag words below this probability (default 0.55)")
    c.add_argument("--no-vad", action="store_true",
                   help="disable voice-activity filtering")
    c.add_argument("--vocab", help="comma-separated terms to bias the recogniser")
    c.add_argument("--vocab-file",
                   help="file of terms — e.g. built from Layer 3's entity_inventory")
    c.add_argument("--terms-from", metavar="MANIFEST",
                   help="Layer 3 manifest — its entity_inventory and deck terms "
                        "drive notation detection, and the vocabulary bias when "
                        "--vocab is not given")
    c.add_argument("--model-dir", help="local model cache (offline hosts)")
    c.set_defaults(fn=cmd_transcript_build)

    c = gs.add_parser("export", help="transcript -> json | html | vtt | txt")
    c.add_argument("file")
    c.add_argument("--format", "-f", default="html",
                   choices=["json", "html", "vtt", "txt"])
    c.add_argument("--out", "-o", required=True)
    c.add_argument("--audio", help="embed this audio so timestamps play (html only)")
    c.add_argument("--template", "-t", help="override transcript-review")
    c.set_defaults(fn=cmd_transcript_export)

    c = gs.add_parser("reflag",
                      help="re-apply notation detection without re-transcribing")
    c.add_argument("file")
    c.add_argument("--terms-from", metavar="MANIFEST", help="Layer 3 manifest")
    c.add_argument("--terms", help="extra comma-separated terms (default: the "
                                   "transcript's own vocab_bias)")
    c.add_argument("--out", "-o", required=True)
    c.set_defaults(fn=cmd_transcript_reflag)

    c = gs.add_parser("apply",
                      help="apply a human review file -> the final transcript")
    c.add_argument("file", help="the transcript to correct")
    c.add_argument("--review", "-r", required=True,
                   metavar="REVIEW.txt", help="the review file (MM:SS | kind | comment)")
    c.add_argument("--by", required=True, metavar="NAME",
                   help="who listened — recorded as confirmed_by on every fix")
    c.add_argument("--keep-going", action="store_true",
                   help="write the transcript even if some review lines fail")
    c.add_argument("--out", "-o", required=True)
    c.set_defaults(fn=cmd_transcript_apply)

    c = gs.add_parser("check", help="validate a transcript's timings")
    c.add_argument("file"); c.set_defaults(fn=cmd_transcript_check)

    # ---- audio
    g = groups.add_parser("audio", help="audio extraction and measurement")
    gs = g.add_subparsers(dest="cmd", metavar="<command>")
    c = gs.add_parser("extract", help="video -> 48 kHz stereo master (RC-003)")
    c.add_argument("file"); c.add_argument("--out", "-o", required=True)
    c.set_defaults(fn=cmd_audio_extract)
    c = gs.add_parser("asr", help="video -> 16 kHz mono copy for alignment")
    c.add_argument("file"); c.add_argument("--out", "-o", required=True)
    c.set_defaults(fn=cmd_audio_asr)
    c = gs.add_parser("probe", help="loudness and format measurements (Layer 7)")
    c.add_argument("file"); c.add_argument("--out", "-o")
    c.set_defaults(fn=cmd_audio_probe)

    # ---- video
    g = groups.add_parser("video", help="video measurement")
    gs = g.add_subparsers(dest="cmd", metavar="<command>")
    c = gs.add_parser("probe", help="resolution/fps/codec/bitrate (TGT-001…008)")
    c.add_argument("file"); c.add_argument("--out", "-o")
    c.set_defaults(fn=cmd_video_probe)
    c = gs.add_parser("slidechanges",
                      help="what is on screen when — changes, slide identity, focus arrow")
    c.add_argument("file")
    c.add_argument("--out", "-o", required=True)
    c.add_argument("--html", help="also write the review page here")
    c.add_argument("--video-id")
    c.add_argument("--deck", help="the .pptx — rendered and matched, so windows carry a "
                                  "slide_id instead of an assumption")
    c.add_argument("--deck-id", help="id prefix for matched slides (default: deck stem)")
    c.add_argument("--fps", type=float, default=4.0)
    c.add_argument("--width", type=int, default=192)
    c.add_argument("--height", type=int, default=108)
    c.add_argument("--threshold", type=float, default=15.0,
                   help="mean-diff at or above this is tiered 'strong' — a measure of how "
                        "much moved, NOT of whether it is a slide turn (default 15)")
    c.add_argument("--weak-threshold", type=float, default=10.0)
    c.add_argument("--frac-threshold", type=float, default=0.8,
                   help="%% of pixels changed that counts as an event — catches "
                        "cross-dissolves the mean cannot see (default 0.8)")
    c.add_argument("--pixel-delta", type=int, default=40,
                   help="a pixel counts as changed above this many levels (default 40)")
    c.add_argument("--min-gap", type=float, default=2.0)
    c.add_argument("--min-window", type=float, default=1.0,
                   help="ignore windows shorter than this (default 1.0s)")
    c.add_argument("--merge-min", type=float, default=0.93,
                   help="two adjacent unmatched windows this alike are one window "
                        "(default 0.93)")
    c.add_argument("--match-min", type=float, default=0.80,
                   help="correlation below this means no deck slide matches (default 0.80)")
    c.add_argument("--match-width", type=int, default=96)
    c.add_argument("--match-height", type=int, default=54)
    c.add_argument("--render-dpi", type=int, default=100)
    c.add_argument("--no-arrow", dest="arrow", action="store_false",
                   help="skip focus-arrow extraction")
    c.add_argument("--arrow-fps", type=float, default=2.0)
    c.add_argument("--arrow-r", type=int, default=170)
    c.add_argument("--arrow-min-px", type=int, default=6)
    c.add_argument("--arrow-min-s", type=float, default=1.0)
    c.add_argument("--no-thumbs", dest="thumbs", action="store_false",
                   help="skip the embedded before/after frames")
    c.add_argument("--thumb-width", type=int, default=560)
    c.add_argument("--thumb-lead", type=float, default=0.6)
    c.set_defaults(fn=cmd_video_slidechanges, arrow=True, thumbs=True)

    c = gs.add_parser("uniquefps", help="unique frames per second (TGT-013)")
    c.add_argument("file"); c.add_argument("--out", "-o")
    c.add_argument("--threshold", type=float, default=24.0)
    c.add_argument("--hi", type=int, default=768); c.add_argument("--lo", type=int, default=320)
    c.set_defaults(fn=cmd_video_uniquefps)

    # ---- check
    g = groups.add_parser("check", help="validate artifacts")
    gs = g.add_subparsers(dest="cmd", metavar="<command>")
    c = gs.add_parser("manifest", help="validate a Layer 3 manifest")
    c.add_argument("file"); c.set_defaults(fn=cmd_check_manifest)

    c = gs.add_parser("representation", help="validate a Layer 5 representation")
    c.add_argument("file")
    c.add_argument("--manifest", help="Layer 3 manifest — enables the element "
                                      "coverage check (VGR-07)")
    c.add_argument("--transcript", help="Layer 4 transcript — reserved")
    c.add_argument("--windows", help="mpk video slidechanges output — enables the "
                                     "slide-agreement and focus-arrow checks")
    c.add_argument("--tolerance", type=float, default=0.05,
                   help="seconds of slack allowed between focus windows (default 0.05)")
    c.set_defaults(fn=cmd_check_representation)

    return ap


def main(argv=None):
    ap = build_parser()
    args = ap.parse_args(argv)
    for attr in ("file", "into", "like"):
        v = getattr(args, attr, None)
        if isinstance(v, str) and not Path(v).exists():
            die(f"{attr}: no such file: {v}")
    if not getattr(args, "fn", None):
        (ap if not args.group else
         ap.parse_args([args.group, "--help"]))
        ap.print_help()
        return 2
    return args.fn(args)


if __name__ == "__main__":
    try:
        sys.exit(main() or 0)
    except BrokenPipeError:
        # `mpk deck info deck.pptx | head` is normal usage; SIGPIPE is not an error.
        try:
            sys.stdout.close()
        finally:
            sys.exit(0)
    except KeyboardInterrupt:
        sys.exit(130)
